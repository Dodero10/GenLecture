import json
import os
import shutil

from pdf2image import convert_from_path
from pptx import Presentation


def create_folder(path):
    """Tạo thư mục nếu chưa tồn tại."""
    if not os.path.exists(path):
        os.makedirs(path)

def extract_metadata_pptx(file_path, output_folder):
    """Trích xuất metadata từ file PPTX."""
    presentation = Presentation(file_path)
    folder_name = os.path.splitext(os.path.basename(file_path))[0]
    slide_folder = os.path.join(output_folder, folder_name)
    metadata = []

    create_folder(slide_folder)

    for i, slide in enumerate(presentation.slides):
        slide_data = {
            "slide_number": i + 1,
            "title": slide.shapes.title.text if slide.shapes.title else "Untitled",
            "text_content": " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")]),
            "images": [],
            "tables": []
        }

        # Extract images
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture shape
                image_name = f"slide_{i + 1}_image_{len(slide_data['images']) + 1}.png"
                image_path = os.path.join(slide_folder, image_name)
                with open(image_path, "wb") as img_file:
                    img_file.write(shape.image.blob)
                slide_data["images"].append(image_name)

        # Extract tables
        for shape in slide.shapes:
            if shape.has_table:
                table_data = []
                table = shape.table
                for row in table.rows:
                    table_data.append([cell.text for cell in row.cells])
                slide_data["tables"].append(table_data)

        metadata.append(slide_data)

    # Save metadata to JSON
    metadata_path = os.path.join(output_folder, f"{folder_name}.json")
    with open(metadata_path, "w", encoding="utf-8") as json_file:
        json.dump(metadata, json_file, ensure_ascii=False, indent=4)

    print(f"Metadata and assets for {file_path} saved to {slide_folder}.")

def extract_metadata_pdf(file_path, output_folder):
    """Trích xuất metadata từ file PDF."""
    folder_name = os.path.splitext(os.path.basename(file_path))[0]
    slide_folder = os.path.join(output_folder, folder_name)
    metadata = []

    create_folder(slide_folder)

    # Convert PDF to images
    pages = convert_from_path(file_path, dpi=300)

    for i, page in enumerate(pages):
        image_name = f"page_{i + 1}.png"
        image_path = os.path.join(slide_folder, image_name)
        page.save(image_path, "PNG")

        slide_data = {
            "page_number": i + 1,
            "images": [image_name],
            "text_content": ""  # Placeholder for OCR or additional text extraction if needed
        }
        metadata.append(slide_data)

    # Save metadata to JSON
    metadata_path = os.path.join(output_folder, f"{folder_name}.json")
    with open(metadata_path, "w", encoding="utf-8") as json_file:
        json.dump(metadata, json_file, ensure_ascii=False, indent=4)

    print(f"Metadata and assets for {file_path} saved to {slide_folder}.")

def extract_metadata(file_path, output_folder="data/metadata"):
    """Xác định loại file và gọi hàm trích xuất phù hợp."""
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    file_extension = os.path.splitext(file_path)[1].lower()
    if file_extension == ".pptx":
        extract_metadata_pptx(file_path, output_folder)
    elif file_extension == ".pdf":
        extract_metadata_pdf(file_path, output_folder)
    else:
        print(f"Unsupported file type: {file_extension}")

# Test example
if __name__ == "__main__":
    test_file_pptx = "example.pptx"
    test_file_pdf = "example.pdf"

    extract_metadata(test_file_pptx)
    extract_metadata(test_file_pdf)
